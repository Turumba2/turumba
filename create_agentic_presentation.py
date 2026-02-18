#!/usr/bin/env python3
"""Generate Agentic AI Workflow Presentation (10-minute lightning talk)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors (same as Turumba 2.0 Overview) ──
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)  # Deep navy
ACCENT_BLUE  = RGBColor(0x38, 0x9C, 0xF7)  # Bright blue
ACCENT_TEAL  = RGBColor(0x06, 0xB6, 0xD4)  # Teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x99, 0x99, 0x99)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
RED          = RGBColor(0xEF, 0x44, 0x44)
PURPLE       = RGBColor(0xA7, 0x8B, 0xFA)
SECTION_BG   = RGBColor(0x14, 0x1F, 0x38)  # Slightly lighter navy

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper Functions (reused from create_presentation.py) ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox

def add_multiline(slide, lines, left, top, width, height, font_size=16, color=WHITE, line_spacing=1.5, bullet=False, font_name="Calibri"):
    """lines: list of (text, bold, color) or just strings"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, bld, clr = line, False, color
        else:
            txt, bld, clr = line[0], line[1] if len(line) > 1 else False, line[2] if len(line) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "\u2022  " if bullet else ""
        p.text = prefix + txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txbox

def add_card(slide, left, top, width, height, title, items, accent=ACCENT_BLUE, title_size=16, item_size=13):
    card = add_shape(slide, left, top, width, height, SECTION_BG, accent, Pt(1.5))
    # Title
    add_text(slide, title, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
             font_size=title_size, color=accent, bold=True)
    # Separator line
    add_rect(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5), Pt(1.5), accent)
    # Items
    if items:
        add_multiline(slide, items, left + Inches(0.25), top + Inches(0.65), width - Inches(0.5),
                       height - Inches(0.85), font_size=item_size, bullet=True, line_spacing=1.4)

def slide_title(slide, title, subtitle=""):
    set_slide_bg(slide, DARK_BG)
    # Top accent line
    add_rect(slide, Inches(0.6), Inches(0.5), Inches(1.5), Pt(3), ACCENT_BLUE)
    # Title
    add_text(slide, title, Inches(0.6), Inches(0.6), Inches(11), Inches(0.6),
             font_size=30, color=WHITE, bold=True, font_name="Calibri Light")
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.15), Inches(11), Inches(0.4),
                 font_size=16, color=MED_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(4.4), Inches(13.333), Pt(2), RGBColor(0x1E, 0x29, 0x3B))

add_text(slide, "Building a Multi-Service SaaS Platform", Inches(1), Inches(1.2), Inches(11), Inches(0.9),
         font_size=48, color=WHITE, bold=True, font_name="Calibri Light")
add_text(slide, "With an Agentic AI Workflow", Inches(1), Inches(2.1), Inches(11), Inches(0.7),
         font_size=44, color=ACCENT_TEAL, bold=True, font_name="Calibri Light")

add_text(slide, "How I used Claude Code to architect, plan, and ship Turumba 2.0\nwith a small team in record time",
         Inches(1), Inches(3.3), Inches(10), Inches(0.9),
         font_size=20, color=LIGHT_GRAY, font_name="Calibri")

add_text(slide, "Lightning Talk  |  10 Minutes",
         Inches(1), Inches(4.8), Inches(6), Inches(0.4),
         font_size=16, color=MED_GRAY)

add_text(slide, "February 2026", Inches(10), Inches(6.5), Inches(2.5), Inches(0.3),
         font_size=14, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: The Challenge
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "The Challenge", "What we set out to build")

# 4 service cards
services = [
    ("Account API", "FastAPI  |  Python 3.11\nUsers, Accounts, Roles,\nContacts, Auth (Cognito)", ACCENT_BLUE),
    ("Messaging API", "FastAPI  |  Python 3.12\nChannels, Messages, Templates,\nGroups, Schedules, Outbox", GREEN),
    ("API Gateway", "KrakenD 2.12.1\nGo plugin, 51 endpoints,\nContext enrichment", ORANGE),
    ("Web Core", "Next.js 16  |  TypeScript\nTurborepo monorepo,\nReact 19, Tailwind v4", PURPLE),
]

for i, (name, desc, color) in enumerate(services):
    x = Inches(0.5) + Inches(i * 3.15)
    y = Inches(1.8)
    card = add_shape(slide, x, y, Inches(2.9), Inches(2.2), SECTION_BG, color, Pt(2))
    add_text(slide, name, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.35),
             font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.5), Pt(1.5), color)
    add_text(slide, desc, x + Inches(0.2), y + Inches(0.75), Inches(2.5), Inches(1.2),
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom challenge statement
add_shape(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6), SECTION_BG, ACCENT_TEAL, Pt(1.5))

challenge_items = [
    ("Multi-tenant, multi-channel platform", True, WHITE),
    ("SMS, Telegram, WhatsApp, Email, SMPP, Messenger", False, LIGHT_GRAY),
    ("Team of 3-4 developers, tight timeline", False, LIGHT_GRAY),
    ("4 separate repos  \u2014  4 services  \u2014  2 databases  \u2014  1 message broker", False, LIGHT_GRAY),
]
add_multiline(slide, challenge_items, Inches(0.9), Inches(4.6), Inches(6.5), Inches(1.8),
              font_size=16, line_spacing=1.6)

add_text(slide, '"Normally requires weeks of upfront design\nand hundreds of tickets"',
         Inches(8), Inches(4.8), Inches(4.5), Inches(1.5),
         font_size=20, color=ORANGE, bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: Agentic vs. Just Using AI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, 'What Makes This "Agentic" and Not Just "Using AI"',
            "The distinction that changes everything")

# Left column: "Using AI"
add_shape(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0), SECTION_BG, RED, Pt(2))
add_text(slide, 'Using AI', Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.5),
         font_size=24, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0.8), Inches(2.5), Inches(5.2), Pt(1.5), RED)
add_multiline(slide, [
    'Ask ChatGPT "how do I build a messaging API?"',
    "Get a generic tutorial response",
    "Copy-paste code snippets",
    "No awareness of your codebase",
    "No awareness of your patterns",
    "Output is conversation, not artifacts",
    "Every session starts from scratch",
], Inches(0.9), Inches(2.8), Inches(5), Inches(3.5),
    font_size=15, bullet=True, line_spacing=1.5, color=LIGHT_GRAY)

# VS divider
add_text(slide, "VS", Inches(6.1), Inches(3.8), Inches(1.1), Inches(0.6),
         font_size=28, color=MED_GRAY, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

# Right column: "Agentic AI"
add_shape(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(5.0), SECTION_BG, GREEN, Pt(2))
add_text(slide, "Agentic AI Workflow", Inches(7), Inches(1.9), Inches(5.8), Inches(0.5),
         font_size=24, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(7.3), Inches(2.5), Inches(5.2), Pt(1.5), GREEN)
add_multiline(slide, [
    ("Context Persistence", True, ACCENT_TEAL),
    "  Reads CLAUDE.md  \u2014  knows your full architecture",
    ("Multi-File Awareness", True, ACCENT_TEAL),
    "  References existing models, patterns, and conventions",
    ("Artifact Production", True, ACCENT_TEAL),
    "  Outputs files: specs, docs, status reports",
    ("Iterative Refinement", True, ACCENT_TEAL),
    "  Collaborator, not oracle  \u2014  adapts to corrections",
], Inches(7.3), Inches(2.8), Inches(5.2), Inches(3.5),
    font_size=14, line_spacing=1.35, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: The Command Center
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "The Command Center", "CLAUDE.md + architecture docs = the agent's knowledge base")

# CLAUDE.md card
add_card(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(3.0),
         "CLAUDE.md  \u2014  System Prompt for the Agent", [
             "300+ lines of precise architectural context",
             "Full gateway routing, context enrichment, backend patterns",
             "CRUDController multi-tenancy behaviors (non-obvious!)",
             "Common commands for every service (run, test, lint, migrate)",
             "Code quality standards and CI/CD configuration",
             "Database models and naming conventions",
         ], ACCENT_BLUE, title_size=15, item_size=12)

# Architecture docs card
add_card(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(3.0),
         "Architecture-First Documentation", [
             ("WHAT_IS_TURUMBA.md  \u2014  Full product specification", True, ACCENT_TEAL),
             ("TURUMBA_MESSAGING.md  \u2014  Messages, templates, events", True, ACCENT_TEAL),
             ("TURUMBA_DELIVERY_CHANNELS.md  \u2014  Channels, credentials", True, ACCENT_TEAL),
             "",
             "Docs serve dual purpose:",
             "  Product spec for humans + context for the agent",
         ], ACCENT_TEAL, title_size=15, item_size=12)

# Key insight box
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), SECTION_BG, ORANGE, Pt(2))
add_text(slide, "The Central Codebase Pattern", Inches(0.8), Inches(5.35), Inches(5), Inches(0.4),
         font_size=18, color=ORANGE, bold=True)
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Pt(1.5), ORANGE)

# Directory structure visualization
add_multiline(slide, [
    ("codebase/", True, WHITE),
    ("  \u251C\u2500 CLAUDE.md                     \u2190  Agent reads this automatically", False, ACCENT_BLUE),
    ("  \u251C\u2500 docs/                           \u2190  Architecture specs, task specs", False, ACCENT_TEAL),
    ("  \u251C\u2500 turumba_account_api/     \u2190  Service repo #1", False, LIGHT_GRAY),
    ("  \u251C\u2500 turumba_messaging_api/  \u2190  Service repo #2", False, LIGHT_GRAY),
    ("  \u251C\u2500 turumba_gateway/            \u2190  Service repo #3", False, LIGHT_GRAY),
    ("  \u2514\u2500 turumba_web_core/          \u2190  Service repo #4", False, LIGHT_GRAY),
], Inches(0.8), Inches(5.85), Inches(8), Inches(1.3), font_size=11, line_spacing=1.15, font_name="Consolas")

add_text(slide, '"The agent does not hallucinate\nfeatures because the features are\nprecisely defined in documents\nit can read."',
         Inches(9.2), Inches(5.85), Inches(3.5), Inches(1.2),
         font_size=14, color=ORANGE, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: AI-Generated Task Specifications
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "AI-Generated Task Specifications",
            "16 task specs (6 BE + 10 FE) \u2014 each precise enough for zero clarifying questions")

# Task spec format card
add_shape(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(5.0), SECTION_BG, ACCENT_BLUE, Pt(2))
add_text(slide, "Task Spec Format", Inches(0.8), Inches(1.95), Inches(5), Inches(0.4),
         font_size=18, color=ACCENT_BLUE, bold=True)
add_rect(slide, Inches(0.8), Inches(2.4), Inches(4.9), Pt(1.5), ACCENT_BLUE)

spec_lines = [
    ("Task ID  |  Title  |  Service  |  Assignee", True, WHITE),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", False, MED_GRAY),
    ("Summary  (what and why)", False, ACCENT_TEAL),
    ("Database model  (columns, types, constraints)", False, ACCENT_TEAL),
    ("Schema definitions  (create, update, response)", False, ACCENT_TEAL),
    ("Controller config  (filters, sorts, schema map)", False, ACCENT_TEAL),
    ("Router endpoints  (methods, paths, req/resp)", False, ACCENT_TEAL),
    ("Step-by-step implementation checklist", False, ACCENT_TEAL),
    ("Testing requirements", False, ACCENT_TEAL),
    ("Definition of done", False, ACCENT_TEAL),
]
add_multiline(slide, spec_lines, Inches(0.8), Inches(2.6), Inches(5), Inches(3.8),
              font_size=13, line_spacing=1.4, font_name="Consolas")

# Right side: key points
add_card(slide, Inches(6.3), Inches(1.8), Inches(6.5), Inches(2.2),
         "Why This Works", [
             "Exact SQLAlchemy model with every column, type, and index",
             "Pydantic schemas for creation, update, and response",
             "Filter/sort configuration with whitelisted operations",
             "A developer can implement without asking a single question",
         ], GREEN, title_size=15, item_size=13)

# Quote highlight
add_shape(slide, Inches(6.3), Inches(4.3), Inches(6.5), Inches(2.5), SECTION_BG, ORANGE, Pt(2))
add_text(slide, '"The spec IS the clarification"', Inches(6.5), Inches(4.5), Inches(6), Inches(0.6),
         font_size=24, color=ORANGE, bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

add_multiline(slide, [
    ("16 task specifications", True, WHITE),
    "  6 backend  +  10 frontend",
    "",
    ("Produced in a fraction of the time", True, WHITE),
    "  with far more consistency and detail",
    "  than manual ticket writing",
], Inches(6.8), Inches(5.2), Inches(5.5), Inches(1.4), font_size=13, line_spacing=1.2, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: From Spec to Implementation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "From Spec to Implementation", "5 complete CRUD entities in 4 days")

# Timeline visualization
milestones = [
    ("Feb 8",  "Messaging API core\narchitecture + dual\ndatabase setup", ACCENT_BLUE),
    ("Feb 9",  "Alembic + pre-commit\n+ pytest infrastructure\ncomplete", ACCENT_TEAL),
    ("Feb 11", "BE-001 (Messages) +\nBE-002 (Channels)\nCRUD complete & closed", GREEN),
    ("Feb 12", "BE-003, BE-004, BE-005\n(Templates, Groups,\nSchedules) all closed", ORANGE),
    ("Feb 13", "Full project audit\n51 gateway endpoints\nconfigured", PURPLE),
]

# Timeline bar
add_rect(slide, Inches(0.8), Inches(3.25), Inches(11.7), Pt(3), ACCENT_BLUE)

for i, (date, desc, color) in enumerate(milestones):
    x = Inches(0.5) + Inches(i * 2.55)
    # Dot on timeline
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.95), Inches(3.1), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # Date above
    add_text(slide, date, x + Inches(0.15), Inches(2.2), Inches(2), Inches(0.4),
             font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description below
    add_shape(slide, x + Inches(0.15), Inches(3.7), Inches(2.2), Inches(1.6), SECTION_BG, color, Pt(1.5))
    add_text(slide, desc, x + Inches(0.25), Inches(3.85), Inches(2), Inches(1.3),
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key callout
add_shape(slide, Inches(0.5), Inches(5.7), Inches(7.5), Inches(1.3), SECTION_BG, GREEN, Pt(2))
add_text(slide, "5 CRUD entities  \u2192  4 days  \u2192  80% test coverage", Inches(0.8), Inches(5.85), Inches(7), Inches(0.5),
         font_size=22, color=GREEN, bold=True)
add_text(slide, "Model + Schema + Controller + Service + Router + Tests for each entity",
         Inches(0.8), Inches(6.35), Inches(7), Inches(0.4),
         font_size=14, color=LIGHT_GRAY)

add_shape(slide, Inches(8.3), Inches(5.7), Inches(4.5), Inches(1.3), SECTION_BG, ORANGE, Pt(2))
add_text(slide, '"Implementation became an\nexecution task, not a design task"',
         Inches(8.3), Inches(5.85), Inches(4.5), Inches(1.0),
         font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Automated Code Review
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "Automated Code Review (Claude Code Action)",
            "Every PR reviewed against the project's architecture  \u2014  automatically")

# Two workflow cards
add_card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.0),
         "claude-code-review.yml  \u2014  Auto Review", [
             "Fires on every PR open / update",
             "Claude reads diff + full codebase context",
             "Posts inline comments + structured summary",
             "No human trigger required",
         ], ACCENT_BLUE, title_size=14, item_size=12)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.0),
         "claude.yml  \u2014  Interactive @claude Review", [
             'Comment "@claude please review" on any PR',
             "Targeted review responding to specific request",
             "Follow-up questions, re-reviews after fixes",
             "Verify whether a concern has been addressed",
         ], GREEN, title_size=14, item_size=12)

# The review prompt details
add_shape(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.1), SECTION_BG, PURPLE, Pt(2))
add_text(slide, "The Review Prompt: 120+ Lines of Architecture Context", Inches(0.8), Inches(4.25), Inches(11), Inches(0.4),
         font_size=18, color=PURPLE, bold=True)
add_rect(slide, Inches(0.8), Inches(4.7), Inches(11.7), Pt(1.5), PURPLE)

# 10 critical patterns in two columns
left_patterns = [
    ("1.", "CRUDController Base Class", "All controllers MUST extend it"),
    ("2.", "Multi-Tenant Scoping", "Every query scoped to x-account-ids"),
    ("3.", "Filter/Sort Config", "Whitelist of allowed filters per entity"),
    ("4.", "Schema Conventions", "PATCH: exclude_unset=True, not exclude_none"),
    ("5.", "Async DB Operations", "Never block the event loop"),
]
right_patterns = [
    ("6.", "PostgreSQL Models", "sa.Uuid(as_uuid=True), not sa.UUID()"),
    ("7.", "Alembic Migrations", "Column types must match models exactly"),
    ("8.", "Testing Standards", "80% coverage, shared conftest fixtures"),
    ("9.", "Domain-Specific Rules", "Status lifecycles, credential handling"),
    ("10.", "Code Quality", "Ruff, proper error chaining, no hardcoded config"),
]

for i, (num, title, desc) in enumerate(left_patterns):
    y = Inches(4.9) + Inches(i * 0.42)
    add_text(slide, num, Inches(0.8), y, Inches(0.3), Inches(0.3), font_size=11, color=PURPLE, bold=True)
    add_text(slide, title, Inches(1.15), y, Inches(2.2), Inches(0.3), font_size=11, color=WHITE, bold=True)
    add_text(slide, desc, Inches(3.4), y, Inches(3), Inches(0.3), font_size=10, color=LIGHT_GRAY)

for i, (num, title, desc) in enumerate(right_patterns):
    y = Inches(4.9) + Inches(i * 0.42)
    add_text(slide, num, Inches(6.8), y, Inches(0.4), Inches(0.3), font_size=11, color=PURPLE, bold=True)
    add_text(slide, title, Inches(7.25), y, Inches(2.2), Inches(0.3), font_size=11, color=WHITE, bold=True)
    add_text(slide, desc, Inches(9.5), y, Inches(3.1), Inches(0.3), font_size=10, color=LIGHT_GRAY)

# Quote at bottom
add_text(slide, '"The prompt IS the reviewer\'s expertise"',
         Inches(3.5), Inches(7.0), Inches(6), Inches(0.4),
         font_size=15, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: Real Security Catches
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "Real Security Catches", "Actual vulnerabilities caught by automated review in production PRs")

# 4 security catch cards
catches = [
    ("SQL Injection in pg_notify.py",
     "PR #24  \u2014  Messaging API",
     "f-string SQL construction allows\narbitrary SQL execution.\nCaught in the event infrastructure PR.",
     RED),
    ("Auth Bug: Cognito Access Tokens",
     "PR #52  \u2014  Account API",
     'AWS Cognito access tokens lack "aud" claim\n\u2014 they use "client_id" instead.\nWould have broken auth for all users.',
     ORANGE),
    ("Multi-Tenant Bypass via Filters",
     "PR #17  \u2014  Messaging API",
     "User-provided account_id filter could\nreplace system scope filter via\n_merge_filters. Cross-tenant data leak.",
     PURPLE),
    ("Delete Skips Account Filtering",
     "PR #51  \u2014  Account API",
     "Deletion by ID did not enforce\naccount_id scoping. Any tenant could\ndelete another tenant's records.",
     ACCENT_BLUE),
]

for i, (title, pr, desc, color) in enumerate(catches):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.8) + Inches(row * 2.55)

    add_shape(slide, x, y, Inches(5.95), Inches(2.3), SECTION_BG, color, Pt(2))

    # CRITICAL badge
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.1), Inches(0.3), color)
    add_text(slide, "CRITICAL", x + Inches(0.15), y + Inches(0.15), Inches(1.1), Inches(0.3),
             font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, title, x + Inches(1.4), y + Inches(0.15), Inches(4.3), Inches(0.35),
             font_size=15, color=color, bold=True)
    add_text(slide, pr, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.25),
             font_size=11, color=MED_GRAY)
    add_rect(slide, x + Inches(0.2), y + Inches(0.85), Inches(5.5), Pt(1), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, desc, x + Inches(0.2), y + Inches(1.0), Inches(5.5), Inches(1.1),
             font_size=12, color=LIGHT_GRAY)

# Bottom insight
add_shape(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45), RGBColor(0x1A, 0x25, 0x38), GREEN, Pt(1))
add_text(slide, "These are not theoretical \u2014 every catch above was in a real PR heading toward production",
         Inches(0.8), Inches(6.95), Inches(11.8), Inches(0.35),
         font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9: The Multi-Round Review Loop
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "The Multi-Round Review Loop",
            "7 rounds on one PR  \u2014  each round goes deeper")

# Flow visualization: horizontal pipeline
flow_steps = [
    ("PR Opened", "Automated review\nfires immediately", ACCENT_BLUE),
    ("Dev Fixes", "Developer addresses\nfirst-wave issues", GREEN),
    ("@claude review", "Tech lead triggers\ntargeted re-review", ORANGE),
    ("Verify Fixes", "Claude confirms\nfixes are correct", ACCENT_TEAL),
    ("Deeper Issues", "Finds new issues\nmissed in round 1", RED),
    ("Dev Fixes Again", "Developer addresses\nnewly found issues", GREEN),
    ("Final Approve", "All critical issues\nresolved \u2192 Merge", PURPLE),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.2) + Inches(i * 1.85)
    y = Inches(1.8)

    add_shape(slide, x, y, Inches(1.6), Inches(1.8), SECTION_BG, color, Pt(1.5))

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y + Inches(0.1), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, str(i + 1), x + Inches(0.55), y + Inches(0.12), Inches(0.5), Inches(0.45),
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, title, x + Inches(0.05), y + Inches(0.7), Inches(1.5), Inches(0.35),
             font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.05), y + Inches(1.1), Inches(1.5), Inches(0.6),
             font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 6:
        add_text(slide, "\u25B6", x + Inches(1.6), y + Inches(0.65), Inches(0.25), Inches(0.4),
                 font_size=14, color=MED_GRAY)

# PR #51 example
add_card(slide, Inches(0.5), Inches(4.0), Inches(8), Inches(3.2),
         "Case Study: Account API PR #51  \u2014  7 Review Rounds", [
             ("Round 1:", True, ACCENT_BLUE),
             "  Cursor Bugbot finds multi-tenant bypass in person-contact retrieval",
             ("Round 2:", True, ACCENT_BLUE),
             "  Claude auto-review finds missing type hints",
             ("Round 3:", True, GREEN),
             '  Dev fixes both. Tech lead: "@claude please review"',
             ("Round 4:", True, ORANGE),
             "  Claude verifies fixes, finds NEW issue: delete operation skips account filtering",
             ("Round 5-7:", True, PURPLE),
             "  Fix \u2192 re-review \u2192 confirm \u2192 APPROVED",
         ], ACCENT_BLUE, title_size=14, item_size=11)

# Pattern highlight
add_shape(slide, Inches(8.8), Inches(4.0), Inches(4), Inches(3.2), SECTION_BG, ORANGE, Pt(2))
add_text(slide, "The Pattern", Inches(9.1), Inches(4.2), Inches(3.4), Inches(0.4),
         font_size=18, color=ORANGE, bold=True)
add_rect(slide, Inches(9.1), Inches(4.65), Inches(3.4), Pt(1.5), ORANGE)
add_multiline(slide, [
    ("1.", True, WHITE),
    "Let automated review catch the first wave",
    "",
    ("2.", True, WHITE),
    "Developer fixes obvious issues",
    "",
    ("3.", True, WHITE),
    '@claude please review to verify + go deeper',
    "",
    ("4.", True, WHITE),
    "Repeat until clean",
], Inches(9.1), Inches(4.85), Inches(3.4), Inches(2.2), font_size=11, line_spacing=1.1, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10: The Team Dynamic
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "The Team Dynamic", "How roles work with an agentic AI workflow")

# Role cards
roles = [
    ("Tech Lead + Claude Code", "Design architecture\nWrite documentation\nGenerate task specs\nAudit progress\nMake merge decisions",
     ACCENT_BLUE, "DESIGN & DECIDE"),
    ("Backend Developer", "Pick up BE specs (BE-001 to BE-006)\nCreate branch, implement spec\nOpen PR, respond to review\nFix issues, request re-review",
     GREEN, "EXECUTE BACKEND"),
    ("Frontend Developer", "Pick up FE specs (FE-001 to FE-010)\nSame precision as BE specs\nImplement UI with exact schemas\nConnected to gateway APIs",
     ORANGE, "EXECUTE FRONTEND"),
    ("Claude Code Action", "Auto-review every PR on open\nCheck all 10 critical patterns\nPost inline comments + summary\nRe-review on @claude mention",
     PURPLE, "REVIEW & VERIFY"),
]

for i, (title, desc, color, badge_text) in enumerate(roles):
    x = Inches(0.3) + Inches(i * 3.25)
    y = Inches(1.8)
    add_shape(slide, x, y, Inches(3), Inches(3.5), SECTION_BG, color, Pt(2))

    # Badge
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.7), Inches(0.35), color)
    add_text(slide, badge_text, x + Inches(0.15), y + Inches(0.15), Inches(2.7), Inches(0.35),
             font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, title, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(0.4),
             font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.6), Pt(1), color)
    add_text(slide, desc, x + Inches(0.2), y + Inches(1.2), Inches(2.6), Inches(2.0),
             font_size=11, color=LIGHT_GRAY)

# Bottom quotes
add_shape(slide, Inches(0.3), Inches(5.6), Inches(6.2), Inches(1.6), SECTION_BG, ACCENT_TEAL, Pt(2))
add_text(slide, '"Developers never waited for a spec.\nThe specs were ready before they\nfinished the previous task."',
         Inches(0.6), Inches(5.8), Inches(5.8), Inches(1.2),
         font_size=17, color=ACCENT_TEAL, bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(6.8), Inches(5.6), Inches(6.2), Inches(1.6), SECTION_BG, ORANGE, Pt(2))
add_text(slide, '"The most valuable thing a tech lead\nproduces is not code  \u2014  it is clarity."',
         Inches(7.1), Inches(5.8), Inches(5.6), Inches(1.2),
         font_size=17, color=ORANGE, bold=True, font_name="Calibri Light", alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11: The Numbers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(slide, "The Numbers", "Metrics from the Turumba 2.0 agentic workflow")

# Big metric cards - row 1
metrics_row1 = [
    ("4", "Services", ACCENT_BLUE),
    ("11", "Backend Entities", GREEN),
    ("51", "Gateway Endpoints", ACCENT_TEAL),
    ("16", "Task Specs Generated", ORANGE),
]

for i, (num, label, color) in enumerate(metrics_row1):
    x = Inches(0.4) + Inches(i * 3.2)
    y = Inches(1.8)
    add_shape(slide, x, y, Inches(2.9), Inches(1.3), SECTION_BG, color, Pt(2))
    add_text(slide, num, x, y + Inches(0.1), Inches(2.9), Inches(0.65),
             font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.8), Inches(2.9), Inches(0.35),
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Big metric cards - row 2
metrics_row2 = [
    ("30+", "PRs Reviewed by Claude", PURPLE),
    ("4 Days", "Zero to Functional API", GREEN),
    ("80%", "Test Coverage Gate", ACCENT_BLUE),
    ("3", "Prompt Iterations", ORANGE),
]

for i, (num, label, color) in enumerate(metrics_row2):
    x = Inches(0.4) + Inches(i * 3.2)
    y = Inches(3.4)
    add_shape(slide, x, y, Inches(2.9), Inches(1.3), SECTION_BG, color, Pt(2))
    add_text(slide, num, x, y + Inches(0.1), Inches(2.9), Inches(0.65),
             font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.8), Inches(2.9), Inches(0.35),
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Critical bugs caught section
add_shape(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.45), SECTION_BG, RED, Pt(1.5))
add_text(slide, "Critical Bugs Caught:   SQL Injection  |  Auth Bypass  |  Tenant Isolation Gaps  |  Delete Operation Bypass",
         Inches(0.7), Inches(5.05), Inches(12), Inches(0.35),
         font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# Prompt evolution
add_shape(slide, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.5), SECTION_BG, ACCENT_BLUE, Pt(1.5))
add_text(slide, "Prompt Evolution: 3 Iterations", Inches(0.7), Inches(5.85), Inches(3), Inches(0.35),
         font_size=16, color=ACCENT_BLUE, bold=True)

evolutions = [
    ("V1: Generic", 'code-review plugin.\nSurface-level feedback.', RED),
    ("V2: Structured", "Custom prompt + auto-review.\nLacked codebase context.", ORANGE),
    ("V3: Architecture-Aware", "120+ lines of patterns.\nCatches real vulnerabilities.", GREEN),
]

for i, (version, desc, color) in enumerate(evolutions):
    x = Inches(0.7) + Inches(i * 4.1)
    y = Inches(6.25)
    # Arrow between
    if i > 0:
        add_text(slide, "\u25B6", x - Inches(0.35), y + Inches(0.15), Inches(0.3), Inches(0.3),
                 font_size=16, color=MED_GRAY)
    add_shape(slide, x, y, Inches(3.6), Inches(0.8), RGBColor(0x1A, 0x25, 0x38), color, Pt(1))
    add_text(slide, version, x + Inches(0.1), y + Inches(0.05), Inches(1.5), Inches(0.3),
             font_size=12, color=color, bold=True)
    add_text(slide, desc, x + Inches(1.6), y + Inches(0.05), Inches(1.9), Inches(0.7),
             font_size=10, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12: Key Takeaways + Thank You
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Decorative bar
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_text(slide, "Key Takeaways", Inches(0.6), Inches(0.5), Inches(11), Inches(0.6),
         font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_rect(slide, Inches(0.6), Inches(1.1), Inches(2), Pt(3), ACCENT_BLUE)

# 5 lessons
lessons = [
    ("1", "Invest heavily in CLAUDE.md", "Quality in = quality out. 300+ lines of architecture context pays dividends on every interaction.", ACCENT_BLUE),
    ("2", "Documentation is infrastructure, not overhead", "Architecture docs become the foundation for specs, audits, and status reports. Machine-readable docs compound.", ACCENT_TEAL),
    ("3", "Specs should be executable, not descriptive", "Exact column definitions + step-by-step checklist = 2 days writing code, not 2 days figuring out the design.", GREEN),
    ("4", "Review prompts need architecture context", "Generic prompts produce generic feedback. 120 lines of patterns = real vulnerability catches.", ORANGE),
    ("5", "The agent accelerates production 10x \u2014 the thinking is still yours", "Every spec, every decision, every review goes through human judgment. The agent is a collaborator, not an oracle.", PURPLE),
]

for i, (num, title, desc, color) in enumerate(lessons):
    y = Inches(1.4) + Inches(i * 0.95)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, num, Inches(0.6), y + Inches(0.07), Inches(0.5), Inches(0.45),
             font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title + description
    add_text(slide, title, Inches(1.3), y + Inches(0.02), Inches(11), Inches(0.35),
             font_size=17, color=color, bold=True)
    add_text(slide, desc, Inches(1.3), y + Inches(0.4), Inches(11), Inches(0.35),
             font_size=12, color=LIGHT_GRAY)

# Thank you + Q&A section
add_rect(slide, Inches(0.6), Inches(6.15), Inches(12), Pt(2), RGBColor(0x1E, 0x29, 0x3B))

add_text(slide, "Thank You  \u2014  Questions?", Inches(0.6), Inches(6.4), Inches(6), Inches(0.6),
         font_size=28, color=ACCENT_TEAL, bold=True, font_name="Calibri Light")

add_text(slide, "Built with Claude Code by Anthropic",
         Inches(8), Inches(6.55), Inches(5), Inches(0.3),
         font_size=14, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ── Save ──
output = "/Users/bengeos/Projects/Digital Strategy/Turumba 2.0/codebase/Agentic_AI_Workflow.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
print(f"Total slides: {len(prs.slides)}")
