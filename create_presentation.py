#!/usr/bin/env python3
"""Generate Turumba 2.0 Overview Presentation (18 slides, ~20 minutes)
Focus: What is Turumba, microservices, responsibilities, architecture, evolution
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE  = RGBColor(0x38, 0x9C, 0xF7)
ACCENT_TEAL  = RGBColor(0x06, 0xB6, 0xD4)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x99, 0x99, 0x99)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
RED          = RGBColor(0xEF, 0x44, 0x44)
PURPLE       = RGBColor(0xA7, 0x8B, 0xFA)
SECTION_BG   = RGBColor(0x14, 0x1F, 0x38)
DARK_CARD    = RGBColor(0x1A, 0x25, 0x38)
PINK         = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ──

def set_bg(slide, color=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, border=None, bw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw or Pt(1)
    else:
        s.line.fill.background()
    return s

def rr(slide, l, t, w, h, fill, border=None, bw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw or Pt(1)
    else:
        s.line.fill.background()
    return s

def circle(slide, l, t, size, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return tb

def ml(slide, lines, l, t, w, h, sz=14, clr=WHITE, sp=1.4, bullet=False, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str): text, bld, c = line, False, clr
        else: text, bld, c = line[0], line[1] if len(line)>1 else False, line[2] if len(line)>2 else clr
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("\u2022  " if bullet else "") + text
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bld
        p.font.name = font; p.space_after = Pt(sz * (sp - 1))
    return tb

def card(slide, l, t, w, h, title, items, accent=ACCENT_BLUE, tsz=15, isz=12):
    rr(slide, l, t, w, h, SECTION_BG, accent, Pt(1.5))
    txt(slide, title, l+Inches(.2), t+Inches(.12), w-Inches(.4), Inches(.35), sz=tsz, clr=accent, bold=True)
    rect(slide, l+Inches(.2), t+Inches(.48), w-Inches(.4), Pt(1.5), accent)
    if items:
        ml(slide, items, l+Inches(.2), t+Inches(.58), w-Inches(.4), h-Inches(.7), sz=isz, bullet=True, sp=1.35)

def section(slide, num, title, subtitle=""):
    set_bg(slide)
    rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(1.1), SECTION_BG)
    rect(slide, Inches(0), Inches(3.2), Inches(.15), Inches(1.1), ACCENT_BLUE)
    txt(slide, f"0{num}" if num<10 else str(num), Inches(.6), Inches(2.2), Inches(2), Inches(.8),
        sz=48, clr=ACCENT_BLUE, bold=True, font="Calibri Light")
    txt(slide, title, Inches(.6), Inches(3.25), Inches(12), Inches(.7),
        sz=36, clr=WHITE, bold=True, font="Calibri Light")
    if subtitle:
        txt(slide, subtitle, Inches(.6), Inches(4.4), Inches(10), Inches(.5), sz=18, clr=MED_GRAY)

def stitle(slide, title, subtitle=""):
    set_bg(slide)
    rect(slide, Inches(.6), Inches(.5), Inches(1.5), Pt(3), ACCENT_BLUE)
    txt(slide, title, Inches(.6), Inches(.6), Inches(11), Inches(.6),
        sz=30, clr=WHITE, bold=True, font="Calibri Light")
    if subtitle:
        txt(slide, subtitle, Inches(.6), Inches(1.15), Inches(11), Inches(.4), sz=16, clr=MED_GRAY)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); set_bg(s)
rect(s, Inches(0), Inches(0), Inches(.15), Inches(7.5), ACCENT_BLUE)
rect(s, Inches(0), Inches(4.4), Inches(13.333), Pt(2), DARK_CARD)

txt(s, "TURUMBA 2.0", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    sz=60, clr=WHITE, bold=True, font="Calibri Light")
txt(s, "Multi-Channel Message Automation Platform", Inches(1), Inches(2.7), Inches(11), Inches(.6),
    sz=28, clr=ACCENT_TEAL, font="Calibri Light")
txt(s, "Send the right message, to the right person,\nthrough the right channel, at the right time.",
    Inches(1), Inches(4.7), Inches(8), Inches(.8), sz=20, clr=LIGHT_GRAY)
txt(s, "Platform Overview  |  Architecture  |  Microservices  |  Evolution",
    Inches(1), Inches(6.2), Inches(10), Inches(.4), sz=14, clr=MED_GRAY)
txt(s, "February 2026", Inches(10), Inches(6.5), Inches(2.5), Inches(.3),
    sz=14, clr=MED_GRAY, align=PP_ALIGN.RIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: What is Turumba Today?
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "What is Turumba Today?", "A working platform with 50+ API endpoints across 4 microservices")

ml(s, [
    ("Turumba 2.0 is a multi-tenant message automation platform that enables", True, WHITE),
    ("organizations to automate communication with their contacts across", False, LIGHT_GRAY),
    ("SMS, SMPP, Telegram, WhatsApp, Messenger, and Email.", False, LIGHT_GRAY),
], Inches(.6), Inches(1.6), Inches(8), Inches(1), sz=17, sp=1.5)

# Feature cards - what's built
features = [
    ("Accounts & Auth", "Multi-tenant accounts, users,\nroles, RBAC, AWS Cognito JWT", GREEN),
    ("Contacts & Groups", "Flexible metadata, custom\nattributes, tags, segmentation", ACCENT_TEAL),
    ("6 Channel Types", "SMS, SMPP, Telegram, WhatsApp,\nMessenger, Email — write-only creds", ACCENT_BLUE),
    ("Template Messages", "{FIRST_NAME} variables with\n6-source resolution + fallbacks", PURPLE),
    ("Group Messaging", "Bulk send with per-recipient\npersonalization + progress tracking", ORANGE),
    ("Scheduled Messages", "One-time & recurring with\ntimezone awareness, pause/resume", RED),
    ("Event Infrastructure", "EventBus + Transactional Outbox\n+ RabbitMQ (zero event loss)", ACCENT_TEAL),
    ("Frontend Apps", "Turumba dashboard + Negarit\nNext.js 16, 24 shared UI components", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(.5) + Inches(col * 3.15)
    y = Inches(3.0) + Inches(row * 2.15)
    rr(s, x, y, Inches(2.95), Inches(1.9), SECTION_BG, color, Pt(1.5))
    txt(s, title, x+Inches(.2), y+Inches(.15), Inches(2.55), Inches(.3), sz=14, clr=color, bold=True)
    rect(s, x+Inches(.2), y+Inches(.48), Inches(2.55), Pt(1), color)
    txt(s, desc, x+Inches(.2), y+Inches(.6), Inches(2.55), Inches(1), sz=11, clr=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: How Turumba Evolves
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "How Turumba Evolves", "The architecture supports growth without rewrites — each layer builds on what's there")

# Stacked layers (bottom to top)
layers = [
    (Inches(5.8), "CURRENT FOUNDATION", "Accounts, Auth, Contacts, Channels, Messages,\nTemplates, Groups, Schedules, Event Infrastructure", GREEN, "BUILT"),
    (Inches(4.2), "HIGH-SCALE DISPATCH", "Channel adapter framework, per-channel dispatch workers,\nwebhook receivers, Redis rate limiting \u2014 1M+ messages/day", ACCENT_BLUE, "DESIGNED"),
    (Inches(2.6), "CONVERSATIONS & SUPPORT", "Omnichannel inbox, bot-first routing, agent assignment,\nreal-time WebSocket service (turumba_realtime)", ORANGE, "DESIGNED"),
    (Inches(1.4), "AI & ANALYTICS", "Intent classification, smart replies, translation,\nsentiment detection, dashboards & reporting", PURPLE, "PLANNED"),
]

for y, title, desc, color, status in layers:
    # Layer bar
    rr(s, Inches(.5), y, Inches(10), Inches(1.35), SECTION_BG, color, Pt(2))
    txt(s, title, Inches(.8), y+Inches(.12), Inches(4), Inches(.3), sz=15, clr=color, bold=True)
    txt(s, desc, Inches(.8), y+Inches(.45), Inches(7), Inches(.7), sz=11, clr=LIGHT_GRAY)
    # Status badge
    badge_color = GREEN if status == "BUILT" else (ACCENT_BLUE if status == "DESIGNED" else MED_GRAY)
    rr(s, Inches(8.5), y+Inches(.15), Inches(1.7), Inches(.28), badge_color)
    txt(s, status, Inches(8.5), y+Inches(.15), Inches(1.7), Inches(.28),
        sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Right-side callout
card(s, Inches(10.8), Inches(1.4), Inches(2.2), Inches(5.7),
     "Key Principle", [
         "Each layer builds on the one below",
         "No rewrites needed",
         "Same DB models, same event pipeline",
         "New workers plug into existing RabbitMQ topology",
         "New models follow the same CRUD pattern",
     ], ACCENT_TEAL, tsz=13, isz=10)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: Section - Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); section(s, 1, "System Architecture", "Microservices, API gateway, and event-driven design")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: System Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "System Overview", "4 services, single gateway entry point, Docker networking")

# Client
rr(s, Inches(4.8), Inches(1.6), Inches(3.6), Inches(.75), SECTION_BG, ACCENT_TEAL, Pt(2))
txt(s, "Turumba Web Apps  (Next.js 16)", Inches(4.8), Inches(1.68), Inches(3.6), Inches(.55),
    sz=14, clr=ACCENT_TEAL, bold=True, align=PP_ALIGN.CENTER)

txt(s, "\u25BC", Inches(6.2), Inches(2.35), Inches(.9), Inches(.3), sz=20, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# Gateway
rr(s, Inches(3), Inches(2.7), Inches(7.3), Inches(1), SECTION_BG, ACCENT_BLUE, Pt(2.5))
txt(s, "KrakenD API Gateway  (Port 8080)", Inches(3), Inches(2.75), Inches(7.3), Inches(.4),
    sz=16, clr=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Context Enrichment  |  Auth Validation  |  Rate Limiting  |  Header Injection",
    Inches(3), Inches(3.15), Inches(7.3), Inches(.35), sz=10, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# Arrows
txt(s, "\u25BC", Inches(4.5), Inches(3.7), Inches(.9), Inches(.3), sz=20, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "\u25BC", Inches(7.8), Inches(3.7), Inches(.9), Inches(.3), sz=20, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# Account API
rr(s, Inches(1.2), Inches(4.0), Inches(5.2), Inches(1.5), SECTION_BG, GREEN, Pt(2))
txt(s, "Account API  (FastAPI / Python 3.11)", Inches(1.2), Inches(4.05), Inches(5.2), Inches(.35),
    sz=14, clr=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Users  |  Accounts  |  Roles  |  Contacts  |  Persons  |  Auth (Cognito)",
    Inches(1.2), Inches(4.4), Inches(5.2), Inches(.3), sz=10, clr=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "PostgreSQL  +  MongoDB  +  AWS Cognito",
    Inches(1.2), Inches(4.75), Inches(5.2), Inches(.3), sz=10, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "7 Routers  |  18 Service Classes  |  /context/basic endpoint",
    Inches(1.2), Inches(5.05), Inches(5.2), Inches(.3), sz=10, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# Messaging API
rr(s, Inches(6.9), Inches(4.0), Inches(5.2), Inches(1.5), SECTION_BG, ORANGE, Pt(2))
txt(s, "Messaging API  (FastAPI / Python 3.12)", Inches(6.9), Inches(4.05), Inches(5.2), Inches(.35),
    sz=14, clr=ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Channels  |  Messages  |  Templates  |  Group Messages  |  Scheduled Messages",
    Inches(6.9), Inches(4.4), Inches(5.2), Inches(.3), sz=10, clr=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "PostgreSQL  +  RabbitMQ  (Transactional Outbox)",
    Inches(6.9), Inches(4.75), Inches(5.2), Inches(.3), sz=10, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "5 Routers  |  15+ Service Classes  |  Outbox Worker",
    Inches(6.9), Inches(5.05), Inches(5.2), Inches(.3), sz=10, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# Arrows to DB
txt(s, "\u25BC", Inches(2.5), Inches(5.5), Inches(.9), Inches(.3), sz=16, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "\u25BC", Inches(4.8), Inches(5.5), Inches(.9), Inches(.3), sz=16, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "\u25BC", Inches(8), Inches(5.5), Inches(.9), Inches(.3), sz=16, clr=MED_GRAY, align=PP_ALIGN.CENTER)
txt(s, "\u25BC", Inches(10.2), Inches(5.5), Inches(.9), Inches(.3), sz=16, clr=MED_GRAY, align=PP_ALIGN.CENTER)

# DBs
dbs = [
    (Inches(1.8), "PostgreSQL", ACCENT_BLUE),
    (Inches(4.3), "MongoDB", ACCENT_TEAL),
    (Inches(7.3), "PostgreSQL", ACCENT_BLUE),
    (Inches(9.6), "RabbitMQ", ORANGE),
]
for x, name, color in dbs:
    rr(s, x, Inches(5.85), Inches(2), Inches(.6), SECTION_BG, color, Pt(1.5))
    txt(s, name, x, Inches(5.9), Inches(2), Inches(.5), sz=11, clr=color, bold=True, align=PP_ALIGN.CENTER)

# Docker network
rr(s, Inches(1.2), Inches(6.7), Inches(10.9), Inches(.45), DARK_CARD, MED_GRAY, Pt(1))
txt(s, "Docker Network: gateway-network  |  All services internal  |  Only gateway exposed on port 8080",
    Inches(1.2), Inches(6.75), Inches(10.9), Inches(.35), sz=11, clr=MED_GRAY, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: How It All Works Together
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "How It All Works Together", "End-to-end flow from user sign-up to message delivery")

steps = [
    ("1", "User signs up", "Web app calls\n/v1/auth/register", ACCENT_BLUE),
    ("2", "Account created", "Account API creates\nuser in Cognito + DB", GREEN),
    ("3", "User logs in", "Receives JWT tokens\n(access, id, refresh)", ACCENT_TEAL),
    ("4", "Context enrichment", "Gateway calls /context/basic\ninjects x-account-ids", ORANGE),
    ("5", "Manage contacts", "Create contacts, groups\nwith flexible metadata", PURPLE),
    ("6", "Send message", "Select channel + template\nAPI renders variables", ACCENT_BLUE),
    ("7", "Events emitted", "EventBus + OutboxMiddleware\natomic DB transaction", GREEN),
    ("8", "Background processing", "Outbox Worker publishes\nto RabbitMQ consumers", ORANGE),
    ("9", "Status tracked", "All activity recorded\nwith delivery status", RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(.5) + Inches(col * 4.2)
    y = Inches(1.7) + Inches(row * 1.85)

    rr(s, x, y, Inches(3.9), Inches(1.6), SECTION_BG, color, Pt(1.5))
    # Number circle
    circle(s, x+Inches(.15), y+Inches(.15), Inches(.5), color)
    txt(s, num, x+Inches(.15), y+Inches(.18), Inches(.5), Inches(.45),
        sz=18, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, x+Inches(.8), y+Inches(.18), Inches(2.8), Inches(.3), sz=14, clr=color, bold=True)
    txt(s, desc, x+Inches(.8), y+Inches(.55), Inches(2.8), Inches(.8), sz=11, clr=LIGHT_GRAY)

    # Arrow between cards in same row
    if col < 2:
        txt(s, "\u25B6", x+Inches(3.9), y+Inches(.55), Inches(.3), Inches(.4), sz=16, clr=MED_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Section - The 4 Services
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); section(s, 2, "The Microservices", "Each service's role and responsibilities")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: turumba_gateway
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "turumba_gateway", "KrakenD 2.12.1 \u2014 Single entry point for the entire platform")

# Context enrichment flow
flow = [
    ("User Request", "Hits /v1/accounts\nwith JWT Bearer token", ACCENT_BLUE),
    ("Context Call", "Gateway calls\n/v1/context/basic\non Account API", ACCENT_TEAL),
    ("Header Injection", "Extracts account_ids\n+ role_ids, injects\nas trusted headers", GREEN),
    ("Anti-Spoofing", "STRIPS any user-\nprovided x-account-ids\nor x-role-ids", RED),
    ("Forward", "Enriched request\nforwarded to target\nbackend service", ORANGE),
]

txt(s, "Context Enrichment Flow", Inches(.6), Inches(1.6), Inches(6), Inches(.35),
    sz=18, clr=ACCENT_BLUE, bold=True)

for i, (title, desc, color) in enumerate(flow):
    x = Inches(.4) + Inches(i * 2.55)
    y = Inches(2.1)
    rr(s, x, y, Inches(2.3), Inches(2.1), SECTION_BG, color, Pt(1.5))
    circle(s, x+Inches(.85), y+Inches(.15), Inches(.55), color)
    txt(s, str(i+1), x+Inches(.85), y+Inches(.18), Inches(.55), Inches(.5),
        sz=18, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, x+Inches(.15), y+Inches(.8), Inches(2), Inches(.3), sz=12, clr=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc, x+Inches(.15), y+Inches(1.15), Inches(2), Inches(.8), sz=10, clr=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "\u25B6", x+Inches(2.3), y+Inches(.75), Inches(.25), Inches(.35), sz=14, clr=MED_GRAY)

# Capabilities
card(s, Inches(.4), Inches(4.5), Inches(4), Inches(2.7),
     "Configuration", [
         "Template-based: krakend.tmpl + partials",
         "Go plugin: context-enricher.so",
         "Lua scripts for request/response mods",
         "File composition via FC_ENABLE=1",
     ], ACCENT_BLUE, tsz=14, isz=11)

card(s, Inches(4.7), Inches(4.5), Inches(4), Inches(2.7),
     "51 Endpoints", [
         "25 Account API routes (auth, users, accounts, roles, contacts)",
         "25 Messaging API routes (channels, messages, templates, groups, schedules)",
         "1 Context route (/v1/context/basic)",
     ], GREEN, tsz=14, isz=11)

card(s, Inches(9), Inches(4.5), Inches(4), Inches(2.7),
     "Pattern Matching", [
         '"POST /v1/accounts" \u2014 exact match',
         '"* /v1/accounts/*" \u2014 single wildcard',
         '"GET /v1/**" \u2014 double wildcard',
         "Bypass list for public endpoints",
     ], ORANGE, tsz=14, isz=11)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9: turumba_account_api
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "turumba_account_api", "FastAPI / Python 3.11 \u2014 Identity, access, and contact management")

# Current entities
entities = [
    ("Users", "Registration, auth via\nAWS Cognito, JWT RS256", GREEN),
    ("Accounts", "Multi-tenant orgs,\nsub-accounts for teams", ACCENT_BLUE),
    ("Roles", "Account-specific with\nJSON permissions", ORANGE),
    ("Account Users", "M:N user-account-role\nmapping table", PURPLE),
    ("Contacts", "MongoDB, flexible metadata,\ncustom attributes, tags", ACCENT_TEAL),
    ("Persons", "MongoDB, person records\nwith attributes", MED_GRAY),
]

txt(s, "Current Entities", Inches(.6), Inches(1.6), Inches(4), Inches(.3), sz=16, clr=GREEN, bold=True)

for i, (name, desc, color) in enumerate(entities):
    col = i % 3
    row = i // 3
    x = Inches(.5) + Inches(col * 4.15)
    y = Inches(2.0) + Inches(row * 1.5)
    rr(s, x, y, Inches(3.85), Inches(1.25), SECTION_BG, color, Pt(1.5))
    txt(s, name, x+Inches(.2), y+Inches(.1), Inches(2), Inches(.3), sz=14, clr=color, bold=True)
    txt(s, desc, x+Inches(.2), y+Inches(.45), Inches(3.4), Inches(.6), sz=11, clr=LIGHT_GRAY)

# Architecture details
card(s, Inches(.5), Inches(5.1), Inches(3.85), Inches(2.1),
     "Architecture", [
         "7 Routers, 18 Service Classes",
         "PostgreSQL (relational) + MongoDB (documents)",
         "3 service classes per entity: Creation, Retrieval, Update",
         "/context/basic powers gateway enrichment",
     ], GREEN, tsz=13, isz=10)

card(s, Inches(4.65), Inches(5.1), Inches(3.85), Inches(2.1),
     "Auth Stack", [
         "AWS Cognito user pool (JWT RS256)",
         "get_current_user, get_current_user_id",
         "require_role('admin') decorator",
         "Multi-account membership per user",
     ], ACCENT_BLUE, tsz=13, isz=10)

card(s, Inches(8.8), Inches(5.1), Inches(4.2), Inches(2.1),
     "Evolution: Agent Preferences", [
         ("AgentPreference model for conversation routing", True, ORANGE),
         "Available channels, topics, working hours",
         "Languages, max concurrent conversations",
         "Online/offline toggle, auto-accept, notifications",
     ], ORANGE, tsz=13, isz=10)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10: turumba_messaging_api
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "turumba_messaging_api", "FastAPI / Python 3.12 \u2014 The messaging core of the platform")

# Current entities row
current = [
    ("Channels", "6 types, JSONB creds,\nwrite-only security", ACCENT_BLUE),
    ("Messages", "Status lifecycle, direction\ntracking, JSONB metadata", GREEN),
    ("Templates", "{VAR} placeholders,\n6-source resolution", PURPLE),
    ("Group Messages", "Bulk send, progress\ntracking, auto-template", ORANGE),
    ("Scheduled Msgs", "One-time / recurring,\ntimezone aware", ACCENT_TEAL),
    ("Outbox Events", "Transactional outbox,\npg_notify, retry logic", RED),
]

txt(s, "Current: 6 Entities (all CRUD implemented)", Inches(.6), Inches(1.6), Inches(8), Inches(.3),
    sz=15, clr=GREEN, bold=True)

for i, (name, desc, color) in enumerate(current):
    x = Inches(.3) + Inches(i * 2.15)
    y = Inches(2.0)
    rr(s, x, y, Inches(2.0), Inches(1.3), SECTION_BG, color, Pt(1.5))
    txt(s, name, x+Inches(.12), y+Inches(.1), Inches(1.76), Inches(.25), sz=11, clr=color, bold=True)
    txt(s, desc, x+Inches(.12), y+Inches(.4), Inches(1.76), Inches(.7), sz=9, clr=LIGHT_GRAY)

# Evolution entities
txt(s, "Evolution: New Models & Infrastructure", Inches(.6), Inches(3.5), Inches(8), Inches(.3),
    sz=15, clr=ORANGE, bold=True)

evo = [
    ("Conversations", "Omnichannel inbox with status\nlifecycle: open > bot > assigned\n> pending > resolved > closed", ORANGE),
    ("Contact Identifiers", "Cross-platform contact resolution:\nsame customer on WhatsApp AND\nTelegram maps to one contact_id", ACCENT_TEAL),
    ("Canned Responses", 'Quick replies via /shortcode\ntrigger: "/greeting", "/refund"\nwith {{contact_name}} variables', PURPLE),
    ("Bot Rules", "Rule-based routing engine:\nkeyword, time-based, channel,\nfallback \u2014 priority-ordered", RED),
    ("Channel Adapters", "Pluggable per-provider:\nTwilio, Telegram Bot API,\nWhatsApp Cloud API, SMPP", ACCENT_BLUE),
    ("Dispatch Workers", "Per-channel-type consumers:\nmessage.dispatch.sms,\n.telegram, .whatsapp, etc.", GREEN),
]

for i, (name, desc, color) in enumerate(evo):
    x = Inches(.3) + Inches(i * 2.15)
    y = Inches(3.9)
    rr(s, x, y, Inches(2.0), Inches(1.65), SECTION_BG, color, Pt(1.5))
    txt(s, name, x+Inches(.12), y+Inches(.1), Inches(1.76), Inches(.25), sz=11, clr=color, bold=True)
    txt(s, desc, x+Inches(.12), y+Inches(.4), Inches(1.76), Inches(1), sz=9, clr=LIGHT_GRAY)

# Bottom stats
rr(s, Inches(.3), Inches(5.8), Inches(12.7), Inches(.5), DARK_CARD, ACCENT_BLUE, Pt(1))
txt(s, "5 Routers  |  15+ Service Classes  |  PostgreSQL + RabbitMQ  |  Outbox Worker  |  80% test coverage (CI)",
    Inches(.3), Inches(5.85), Inches(12.7), Inches(.4), sz=12, clr=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Message status flow at bottom
txt(s, "Message Lifecycle:", Inches(.5), Inches(6.5), Inches(2), Inches(.3), sz=12, clr=MED_GRAY, bold=True)
statuses = ["Queued", "Sending", "Sent", "Delivered"]
for i, st in enumerate(statuses):
    x = Inches(2.5) + Inches(i * 2)
    color = [ACCENT_BLUE, ORANGE, GREEN, GREEN][i]
    rr(s, x, Inches(6.5), Inches(1.4), Inches(.35), color)
    txt(s, st, x, Inches(6.5), Inches(1.4), Inches(.35), sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "\u25B6", x+Inches(1.4), Inches(6.45), Inches(.5), Inches(.35), sz=12, clr=MED_GRAY, align=PP_ALIGN.CENTER)
# Failed branch
txt(s, "or  Failed \u25B6 Retry", Inches(10.5), Inches(6.5), Inches(2.5), Inches(.35), sz=10, clr=RED, bold=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11: turumba_web_core
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "turumba_web_core", "Turborepo monorepo \u2014 Next.js 16, TypeScript, Tailwind v4")

# Apps
card(s, Inches(.5), Inches(1.7), Inches(4), Inches(2.6),
     "Turumba App  (Port 3600)", [
         "Full-featured message automation dashboard",
         "Account, team, and contact management",
         "Send, schedule, and group messages",
         "Template and channel management",
         "Monitor delivery status and activity",
     ], ACCENT_BLUE, tsz=14, isz=11)

card(s, Inches(4.8), Inches(1.7), Inches(4), Inches(2.6),
     "Negarit App  (Port 3500)", [
         "Streamlined messaging-focused app",
         "Send and receive messages only",
         "Schedule messages for future delivery",
         "Message history and delivery status",
         "Lightweight alternative to full dashboard",
     ], ACCENT_TEAL, tsz=14, isz=11)

card(s, Inches(9.1), Inches(1.7), Inches(3.9), Inches(2.6),
     "Shared Packages", [
         "@repo/ui \u2014 24 Radix-based components",
         "@repo/eslint-config \u2014 shared lint rules",
         "@repo/typescript-config \u2014 shared tsconfig",
         "Field composition system for forms",
         "Tailwind v4 with oklch color tokens",
     ], PURPLE, tsz=14, isz=11)

# Built features
card(s, Inches(.5), Inches(4.6), Inches(6.1), Inches(2.6),
     "What's Built", [
         ("Auth: Sign in, Sign up, Email verification, TOTP 2FA", True, GREEN),
         ("Server-side auth guard (middleware)", False, GREEN),
         ("Organization management (create, switch, settings)", False, GREEN),
         ("User management within organizations", False, GREEN),
         ("Generic Table Builder with pagination", False, GREEN),
         ("AWS Amplify 6.16 + Cognito integration", False, GREEN),
     ], GREEN, tsz=14, isz=11)

card(s, Inches(6.9), Inches(4.6), Inches(6.1), Inches(2.6),
     "Planned: 10 Messaging Pages", [
         "FE-002/03 \u2014 Delivery Channels table + create",
         "FE-005/06 \u2014 Templates table + create/edit",
         "FE-004/01 \u2014 Messages table + new message compose",
         "FE-007/08 \u2014 Group messages table + create",
         "FE-009/10 \u2014 Scheduled messages table + create/edit",
         ("+ Conversation inbox UI (future)", True, ORANGE),
     ], ORANGE, tsz=14, isz=11)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12: Section - Deep Dives
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); section(s, 3, "Under the Hood", "Multi-tenancy and event-driven architecture")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13: Multi-Tenancy + Event Architecture (combined)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "Multi-Tenancy & Event Architecture")

# LEFT: 3-Layer Tenant Isolation
txt(s, "3-Layer Tenant Isolation", Inches(.5), Inches(1.6), Inches(6), Inches(.35),
    sz=16, clr=ACCENT_BLUE, bold=True)

isolation_layers = [
    ("Layer 1: Gateway", "Context-enricher Go plugin resolves user \u2192 account.\nInjects x-account-ids, x-role-ids headers.\nSTRIPS any user-provided values (anti-spoofing).", ACCENT_BLUE),
    ("Layer 2: Controller", "Default filter: account_id:eq:{header_value}.\n\"Trusted system filter\" \u2014 bypasses user validation.\nCannot be overridden by query parameters.", GREEN),
    ("Layer 3: Service", "set_header_context(headers) extracts IDs.\nAll DB queries scoped to injected account.\nRole-based access control per operation.", ORANGE),
]

for i, (title, desc, color) in enumerate(isolation_layers):
    y = Inches(2.05) + Inches(i * 1.7)
    rr(s, Inches(.5), y, Inches(6.1), Inches(1.5), SECTION_BG, color, Pt(1.5))
    txt(s, title, Inches(.7), y+Inches(.1), Inches(3), Inches(.3), sz=13, clr=color, bold=True)
    txt(s, desc, Inches(.7), y+Inches(.45), Inches(5.5), Inches(.9), sz=10, clr=LIGHT_GRAY)

# RIGHT: Transactional Outbox Pipeline
txt(s, "Transactional Outbox Pipeline", Inches(6.9), Inches(1.6), Inches(6), Inches(.35),
    sz=16, clr=ORANGE, bold=True)

pipeline_steps = [
    ("1. EventBus", "Controller emits domain events.\nIn-memory, request-scoped.\nNo persistence yet.", ACCENT_BLUE),
    ("2. Outbox Middleware", "Flushes events to outbox_events\ntable in SAME DB transaction.\nAtomic: entity + events.", ACCENT_TEAL),
    ("3. db.commit()", "Single commit persists both\nentity changes AND outbox events.\nFail = both rolled back.", GREEN),
    ("4. Outbox Worker", "pg_notify wakes worker instantly.\nPublishes to RabbitMQ exchange.\nrouting_key = event_type.", ORANGE),
    ("5. Consumers", "Process events: group message\nexpansion, schedule triggers,\ndispatch to channels.", RED),
]

for i, (title, desc, color) in enumerate(pipeline_steps):
    y = Inches(2.05) + Inches(i * 1.02)
    rr(s, Inches(6.9), y, Inches(6.1), Inches(.88), SECTION_BG, color, Pt(1.5))
    txt(s, title, Inches(7.1), y+Inches(.05), Inches(2.2), Inches(.25), sz=11, clr=color, bold=True)
    txt(s, desc, Inches(9.3), y+Inches(.05), Inches(3.5), Inches(.75), sz=9, clr=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14: Section - Evolution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); section(s, 4, "Platform Evolution", "High-scale dispatch and customer support")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15: High-Scale Messaging Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "High-Scale Messaging Architecture", "Designed for 1M+ messages/day with per-channel scaling")

# Channel Adapter Layer
card(s, Inches(.4), Inches(1.6), Inches(4.1), Inches(2.8),
     "Channel Adapter Layer", [
         "Pluggable adapter per channel type + provider",
         "Common interface: send(), verify_credentials(),\n  check_health(), parse_webhook()",
         "SMS: Twilio, Africa's Talking, Vonage adapters",
         "Telegram Bot API, WhatsApp Cloud API",
         "SMPP: persistent TCP to SMSCs",
         "Messenger (Graph API), Email (SMTP)",
         "Adapter registry: channel_type + provider \u2192 class",
     ], ACCENT_BLUE, tsz=14, isz=10)

# Two-Stage Dispatch Pipeline
card(s, Inches(4.7), Inches(1.6), Inches(4.2), Inches(2.8),
     "Two-Stage Dispatch Pipeline", [
         ("Stage 1: Fan-Out (Group Messages)", True, ORANGE),
         "Fetches contacts in batches of 1,000",
         "Renders template per contact",
         "Batch-inserts Message records (queued)",
         "Publishes N dispatch events to channel queues",
         ("Stage 2: Per-Channel Dispatch", True, GREEN),
         "Each channel type has dedicated RabbitMQ queue",
         "Workers load creds from Redis cache",
         "Call adapter.send(), update status",
     ], ORANGE, tsz=14, isz=10)

# Per-channel queues + rate limiting
card(s, Inches(9.1), Inches(1.6), Inches(3.9), Inches(2.8),
     "Per-Channel Queues", [
         "message.dispatch.sms",
         "message.dispatch.telegram",
         "message.dispatch.whatsapp",
         "message.dispatch.messenger",
         "message.dispatch.email",
         "message.dispatch.smpp",
         "message.status.update",
         "webhook.inbound",
         ("Independent scaling per channel", True, ACCENT_TEAL),
     ], ACCENT_TEAL, tsz=14, isz=10)

# Bottom row
card(s, Inches(.4), Inches(4.7), Inches(4.1), Inches(2.5),
     "Webhook Receivers", [
         "Inbound messages + delivery status from providers",
         "Verify HMAC signature per provider",
         "Return 200 immediately (< 1 second)",
         "Enqueue to RabbitMQ for async processing",
         "Idempotent: deduplicate by provider message ID",
         "Each provider has different HMAC scheme",
     ], RED, tsz=14, isz=10)

card(s, Inches(4.7), Inches(4.7), Inches(4.2), Inches(2.5),
     "Rate Limiting (3 Levels)", [
         ("Per-channel instance", True, ACCENT_BLUE),
         "  Redis token bucket (channel.rate_limit)",
         ("Per-provider global", True, ORANGE),
         "  Account-level limits (e.g., Twilio 100 msg/sec)",
         ("Per-tenant quota", True, PURPLE),
         "  Daily/monthly caps (free vs. pro tier)",
     ], PURPLE, tsz=14, isz=10)

card(s, Inches(9.1), Inches(4.7), Inches(3.9), Inches(2.5),
     "New Infrastructure", [
         ("Redis", True, RED),
         "  Rate limiting, credential cache, progress",
         "  counters, dedup locks, channel health",
         ("SMPP Gateway (Jasmin)", True, ORANGE),
         "  Persistent TCP to SMSCs",
         ("PostgreSQL Read Replica", True, GREEN),
         "  Separate read/write paths at scale",
         ("Table Partitioning", True, ACCENT_BLUE),
         "  Messages partitioned by month",
     ], RED, tsz=14, isz=10)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16: Conversations & Customer Support
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "Conversations & Customer Support", "Omnichannel inbox with bot-first routing and real-time push")

# turumba_realtime (new service)
card(s, Inches(.4), Inches(1.6), Inches(4.1), Inches(2.5),
     "NEW: turumba_realtime", [
         ("5th microservice \u2014 Node.js + Socket.IO", True, PINK),
         "Subscribes to RabbitMQ conversation events",
         "Pushes to connected browsers via WebSocket",
         "Redis adapter for horizontal scaling",
         "Presence tracking + typing indicators",
         "Two namespaces: /agents, /customers",
         "JWT auth on WebSocket handshake",
     ], PINK, tsz=14, isz=10)

# Conversation Model
card(s, Inches(4.7), Inches(1.6), Inches(4.2), Inches(2.5),
     "Conversation Model", [
         ("Status lifecycle:", True, ACCENT_TEAL),
         "  open \u2192 bot \u2192 assigned \u2192 pending \u2192 resolved \u2192 closed",
         "ContactIdentifier: cross-platform resolution",
         "  Same customer on WhatsApp AND Telegram \u2192 one contact",
         "CannedResponses: /greeting, /refund shortcuts",
         "Internal notes (is_private: true) for agents",
         "SLA tracking: first_reply_at, resolved_at",
     ], ACCENT_TEAL, tsz=14, isz=10)

# Bot-First Routing
card(s, Inches(9.1), Inches(1.6), Inches(3.9), Inches(2.5),
     "Bot-First Routing", [
         ("Phase 1: Rule-Based (MVP)", True, GREEN),
         "  Keyword matching, time-based,\n  channel routing, fallback",
         ("Phase 2: AI Intent", True, ORANGE),
         "  LLM classifies intent +\n  confidence threshold",
         ("Phase 3: Conversational Bot", True, PURPLE),
         "  Multi-turn, knowledge base,\n  handoff to human",
     ], GREEN, tsz=14, isz=10)

# Inbound Flow
txt(s, "Inbound Conversation Flow", Inches(.5), Inches(4.3), Inches(6), Inches(.3),
    sz=15, clr=ORANGE, bold=True)

inbound = [
    ("Customer\nmessages", "WhatsApp, Telegram,\nSMS, Messenger...", ACCENT_BLUE),
    ("Webhook\nReceiver", "Verify HMAC, return\n200, enqueue", ACCENT_TEAL),
    ("Inbound\nWorker", "Resolve contact,\nfind/create convo", GREEN),
    ("Bot\nRouter", "Evaluate rules,\nauto-reply, label", ORANGE),
    ("Agent\nRouting", "Filter by availability,\nassign round-robin", PURPLE),
    ("Real-Time\nPush", "Socket.IO pushes\nto agent inbox", PINK),
]

for i, (title, desc, color) in enumerate(inbound):
    x = Inches(.3) + Inches(i * 2.15)
    y = Inches(4.75)
    rr(s, x, y, Inches(1.95), Inches(1.65), SECTION_BG, color, Pt(1.5))
    txt(s, title, x+Inches(.1), y+Inches(.1), Inches(1.75), Inches(.5), sz=11, clr=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc, x+Inches(.1), y+Inches(.65), Inches(1.75), Inches(.7), sz=9, clr=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, "\u25B6", x+Inches(1.95), y+Inches(.55), Inches(.2), Inches(.35), sz=12, clr=MED_GRAY)

# Agent routing algorithm
rr(s, Inches(.3), Inches(6.6), Inches(12.7), Inches(.65), DARK_CARD, ACCENT_TEAL, Pt(1))
txt(s, "Agent Routing:  Filter by is_available + working hours + available_channels + topics + capacity  \u25B6  Sort by least active + longest idle  \u25B6  Assign",
    Inches(.6), Inches(6.65), Inches(12.2), Inches(.55), sz=11, clr=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 17: Technology Stack
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); stitle(s, "Technology Stack", "Current technologies + planned additions")

categories = [
    ("API Gateway", [("KrakenD 2.12.1", "Go plugins, Lua scripts")], ACCENT_BLUE),
    ("Backend", [("FastAPI", "Python 3.11 / 3.12"), ("SQLAlchemy", "Async ORM"), ("Motor", "MongoDB async")], GREEN),
    ("Auth", [("AWS Cognito", "JWT RS256"), ("Amplify 6.16", "Frontend SDK")], ORANGE),
    ("Databases", [("PostgreSQL", "Relational data"), ("MongoDB", "Document data"), ("RabbitMQ", "Message broker")], PURPLE),
    ("Frontend", [("Next.js 16", "App Router"), ("TypeScript", "Strict mode"), ("Tailwind v4", "oklch tokens")], ACCENT_TEAL),
    ("UI / Forms", [("Radix UI", "Accessible primitives"), ("React Hook Form", "+ Zod validation")], RED),
    ("DevOps", [("Docker", "Compose orchestration"), ("GitHub Actions", "CI/CD pipelines"), ("Turborepo", "Monorepo builds")], ACCENT_BLUE),
    ("Planned", [("Redis", "Rate limiting, cache, presence"), ("Socket.IO", "Real-time WebSocket"), ("Jasmin", "SMPP gateway")], PINK),
]

for i, (cat, items, color) in enumerate(categories):
    col = i % 4
    row = i // 4
    x = Inches(.4) + Inches(col * 3.2)
    y = Inches(1.7) + Inches(row * 2.8)
    h = Inches(2.5)
    rr(s, x, y, Inches(2.95), h, SECTION_BG, color, Pt(1.5))
    txt(s, cat, x+Inches(.2), y+Inches(.12), Inches(2.55), Inches(.3), sz=14, clr=color, bold=True)
    rect(s, x+Inches(.2), y+Inches(.45), Inches(2.55), Pt(1), color)
    item_lines = [(f"{t}  \u2014  {d}", False, LIGHT_GRAY) for t, d in items]
    ml(s, item_lines, x+Inches(.2), y+Inches(.55), Inches(2.55), h-Inches(.7), sz=11, bullet=True, sp=1.6)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 18: Thank You / Q&A
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = blank(); set_bg(s)
rect(s, Inches(0), Inches(0), Inches(.15), Inches(7.5), ACCENT_BLUE)

txt(s, "Thank You", Inches(1), Inches(1.8), Inches(11), Inches(1),
    sz=52, clr=WHITE, bold=True, font="Calibri Light")
txt(s, "Questions & Discussion", Inches(1), Inches(3), Inches(11), Inches(.6),
    sz=28, clr=ACCENT_TEAL, font="Calibri Light")

rect(s, Inches(1), Inches(3.9), Inches(4), Pt(2), DARK_CARD)

ml(s, [
    "4 microservices today, evolving to 5 (turumba_realtime)",
    "51 API endpoints  |  12 data entities  |  6 messaging channels",
    "Multi-tenant SaaS with 3-layer security",
    "Event-driven architecture (transactional outbox + RabbitMQ)",
    "Designed to scale to 1M+ messages/day",
    "Conversation inbox with bot-first routing on the roadmap",
], Inches(1), Inches(4.2), Inches(8), Inches(2.5), sz=15, clr=LIGHT_GRAY, sp=1.7, bullet=True)


# ── Save ──
output = "/Users/bengeos/Projects/Digital Strategy/Turumba 2.0/codebase/Turumba_2.0_Overview.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
print(f"Total slides: {len(prs.slides)}")
